from flask import Flask, render_template, request, send_file
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import date
import textwrap, requests, io, os

app = Flask(__name__, template_folder="templates")

# ---------------- CONFIG ----------------
GEMINI_KEY = "AIzaSyC6NPdnA8qHS2EUzghHE7rSHZ-EqD9kW-w"  # put your real key here
PIXABAY_KEY = "53014821-98af160f6010f41992d35e0ac"
genai.configure(api_key=GEMINI_KEY)
MODEL = genai.GenerativeModel("models/gemini-2.0-flash")

# ---------------- HELPERS ----------------
def fetch_image(query):
    try:
        url = f"https://pixabay.com/api/?key={PIXABAY_KEY}&q={requests.utils.quote(query)}&image_type=photo&orientation=horizontal&per_page=3"
        r = requests.get(url, timeout=8)
        hits = r.json().get("hits", [])
        if hits:
            link = hits[0].get("largeImageURL") or hits[0].get("webformatURL")
            if link:
                resp = requests.get(link, timeout=8)
                if resp.status_code == 200:
                    return io.BytesIO(resp.content)
    except:
        pass
    return None


def get_slides(topic):
    prompt = f"""
You are a presentation content creator.

Make exactly 8 PowerPoint slides about "{topic}".
For each slide, follow this format:
Title: <short clear title, 3–6 words>
- <Bullet point 1: 1–2 short sentences (about 20–35 words)>
- <Bullet point 2: 1–2 short sentences>
- <Bullet point 3: 1–2 short sentences>
- <Bullet point 4: 1–2 short sentences>
The 8th slide must be titled "Conclusion" and summarize everything in 4 clear points.
Do not use markdown or numbering, and return only plain text in the above format.
"""
    try:
        resp = MODEL.generate_content(prompt)
        raw = resp.text.strip()
    except Exception:
        raw = ""

    slides, current = [], None
    for line in raw.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.lower().startswith("title:"):
            if current:
                slides.append(current)
            title = line.split(":", 1)[1].strip()
            current = {"title": title, "bullets": []}
        elif line.startswith("-"):
            if current:
                current["bullets"].append(line[1:].strip())
    if current:
        slides.append(current)

    # fallback if AI fails
    if len(slides) < 8:
        slides = [{"title": f"{topic} Slide {i+1}", "bullets": [
            f"Main point {j+1} about {topic}" for j in range(4)
        ]} for i in range(8)]

    return slides[:8]


def create_ppt(topic, slides):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # title slide
    title_slide = prs.slides.add_slide(blank)
    img = fetch_image(topic)
    if img:
        title_slide.shapes.add_picture(img, 0, 0, prs.slide_width, prs.slide_height)

    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10.5), Inches(2))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.add_paragraph()
    p.text = topic
    p.font.size, p.font.bold = Pt(55), True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # content slides
    for s in slides[1:]:
        slide = prs.slides.add_slide(blank)
        bg = fetch_image(s["title"]) or fetch_image(topic)
        if bg:
            slide.shapes.add_picture(bg, 0, 0, prs.slide_width, prs.slide_height)

        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10.7), Inches(1.2))
        tbox.text_frame.word_wrap = True
        tp = tbox.text_frame.add_paragraph()
        tp.text = s["title"]
        tp.font.size, tp.font.bold = Pt(38), True
        tp.font.color.rgb = RGBColor(255, 255, 255)

        cbox = slide.shapes.add_textbox(Inches(1), Inches(1.9), Inches(8.5), Inches(4.7))
        tf = cbox.text_frame
        tf.word_wrap = True
        for b in s["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(26)
            p.font.color.rgb = RGBColor(255, 255, 255)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf



@app.route("/")
def home():
    return render_template("index.html")


@app.route("/generate", methods=["POST"])
def generate():
    topic = request.form.get("topic", "").strip()
    if not topic:
        topic = "AI Presentation"
    slides = get_slides(topic)
    ppt = create_ppt(topic, slides)
    filename = f"{topic.replace(' ', '_')}_AI_Presentation.pptx"
    with open(filename, "wb") as f:
        f.write(ppt.getbuffer())
    return send_file(filename, as_attachment=True)


if __name__ == "__main__":
    from waitress import serve
    import os
    port = int(os.environ.get("PORT", 10000))
    serve(app, host="0.0.0.0", port=port)


